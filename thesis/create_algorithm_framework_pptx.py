"""Create an editable PowerPoint version of the ForesightNav framework.

The architecture and training paths follow commit
6e137e952fc075b3804c9addadc074691482df02. Every visible element is a
native PowerPoint shape, connector, or text box; no flattened figure is used.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "figures" / "fig_foresightnav_framework_editable.pptx"


COLORS = {
    "ink": "20252B",
    "muted": "68737D",
    "line": "3D4852",
    "blue": "3D8FC4",
    "blue_fill": "EAF4FA",
    "orange": "D97A24",
    "orange_fill": "FFF1E3",
    "gold": "C69A18",
    "gold_fill": "FFF8DE",
    "navy": "244F73",
    "navy_fill": "E9F0F6",
    "green": "5B8F68",
    "green_fill": "EDF6EF",
    "teal": "119A8D",
    "teal_fill": "E8F7F4",
    "coral": "C95537",
    "coral_fill": "FCEDE8",
    "purple": "7A4DC2",
    "purple_fill": "FCFAFF",
    "band": "FAFAF8",
    "band_edge": "B5BEC5",
    "white": "FFFFFF",
}


def rgb(value):
    return RGBColor.from_string(COLORS.get(value, value))


def set_shape_style(shape, edge, fill="FFFFFF", width=1.2, dashed=False):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(edge)
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def add_box(slide, x, y, w, h, edge, fill="FFFFFF", dashed=False, radius=True, width=1.2):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    return set_shape_style(shape, edge, fill, width, dashed)


def set_text(
    shape,
    text,
    size=10,
    color="ink",
    bold=False,
    italic=False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    font="Times New Roman",
    margins=(0.03, 0.03, 0.02, 0.02),
):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margins[0])
    tf.margin_right = Inches(margins[1])
    tf.margin_top = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return shape


def add_text(slide, x, y, w, h, text, **kwargs):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(shape, text, **kwargs)


def add_arrowhead(connector, at_end=True):
    ln = connector.line._get_or_add_ln()
    tag = "a:tailEnd" if at_end else "a:headEnd"
    node = OxmlElement(tag)
    node.set("type", "triangle")
    node.set("w", "sm")
    node.set("len", "sm")
    ln.append(node)


def add_line(slide, x1, y1, x2, y2, color="line", width=1.1, dashed=False, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        add_arrowhead(line)
    return line


def add_route(slide, points, color="line", width=1.1, dashed=False, arrow=True):
    lines = []
    for index, ((x1, y1), (x2, y2)) in enumerate(zip(points[:-1], points[1:])):
        lines.append(
            add_line(
                slide,
                x1,
                y1,
                x2,
                y2,
                color=color,
                width=width,
                dashed=dashed,
                arrow=arrow and index == len(points) - 2,
            )
        )
    return lines


def add_heading(slide, x, w, text):
    add_text(slide, x, 0.08, w, 0.28, text, size=13.2, bold=True)


def add_lidar_icon(slide, x, y, color="blue"):
    for row in range(3):
        for col in range(6):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + col * 0.085),
                Inches(y + row * 0.085),
                Inches(0.058),
                Inches(0.058),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(color if (row + col) % 3 == 0 else "white")
            cell.line.color.rgb = rgb(color)
            cell.line.width = Pt(0.5)


def add_encoder(slide, x, y, w, h, edge, fill, title, dim):
    add_box(slide, x, y, w, h, edge, fill, dashed=True, width=1.2)
    sizes = [(0.34, 0.55), (0.29, 0.46), (0.21, 0.36)]
    starts = [x + 0.24, x + 0.78, x + 1.24]
    for sx, (sw, sh) in zip(starts, sizes):
        layer = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(sx),
            Inches(y + (h - sh) / 2 + 0.04),
            Inches(sw),
            Inches(sh),
        )
        layer.rotation = 94
        set_shape_style(layer, "ink", edge, width=0.7)
    add_text(slide, x + 0.25, y + 0.04, w - 0.50, 0.23, title, size=10.1, bold=True)
    add_text(slide, x + w - 0.52, y + h - 0.23, 0.44, 0.18, dim, size=8.7, color=edge, align=PP_ALIGN.RIGHT)


def add_network(slide, x, y, w, h, edge, fill, title, subtitle):
    add_text(slide, x, y - 0.27, w, 0.24, title, size=11.0, bold=True)
    add_box(slide, x, y, w, h, edge, fill, width=1.25)
    counts = [3, 4, 2]
    xs = [x + 0.34, x + w / 2, x + w - 0.34]
    layers = []
    for count, cx in zip(counts, xs):
        ys = [y + 0.25 + idx * (h - 0.50) / max(count - 1, 1) for idx in range(count)]
        layers.append([(cx, cy) for cy in ys])
    for left, right in zip(layers[:-1], layers[1:]):
        for x1, y1 in left:
            for x2, y2 in right:
                add_line(slide, x1, y1, x2, y2, color=edge, width=0.45)
    for layer in layers:
        for cx, cy in layer:
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - 0.055),
                Inches(cy - 0.055),
                Inches(0.11),
                Inches(0.11),
            )
            set_shape_style(node, edge, "white", width=0.8)
    add_text(slide, x, y + h + 0.02, w, 0.19, subtitle, size=8.4, color="muted", italic=True)


def add_vo_icon(slide, x, y):
    cone = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(0.48), Inches(0.40))
    cone.rotation = 90
    set_shape_style(cone, "coral", "coral_fill", width=0.8)
    cap = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.34), Inches(y + 0.13), Inches(0.13), Inches(0.13))
    set_shape_style(cap, "coral", "white", width=0.8)
    add_line(slide, x - 0.12, y + 0.36, x + 0.25, y + 0.21, color="navy", width=0.8, arrow=True)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")

    # Section headings.
    add_heading(slide, 0.18, 2.10, "Local observation")
    add_heading(slide, 2.65, 2.05, "Modality encoders")
    add_heading(slide, 5.35, 2.70, "Shared representation")
    add_heading(slide, 8.15, 3.05, "Predictive actor-critic")
    add_heading(slide, 11.35, 1.75, "Command")

    # Observation block.
    add_box(slide, 0.12, 0.45, 2.22, 3.70, "purple", "purple_fill", dashed=True, width=1.25)
    add_box(slide, 0.30, 0.76, 1.88, 0.78, "blue", "blue_fill", width=1.0)
    add_lidar_icon(slide, 0.43, 0.98)
    add_text(slide, 1.12, 0.92, 0.92, 0.22, "LiDAR Lₜ", size=10.2, bold=True, color="blue")
    add_text(slide, 1.11, 1.17, 0.94, 0.20, "36 × 4 proximity", size=8.0, color="muted")

    add_box(slide, 0.30, 1.87, 1.88, 0.78, "orange", "orange_fill", width=1.0)
    add_text(slide, 0.40, 2.03, 1.68, 0.25, "Dynamic obstacles Dₜ", size=9.4, bold=True, color="orange")
    add_text(slide, 0.40, 2.31, 1.68, 0.19, "5 × 10 ordered states", size=8.0, color="muted")

    add_box(slide, 0.30, 2.98, 1.88, 0.78, "gold", "gold_fill", width=1.0)
    add_text(slide, 0.40, 3.14, 1.68, 0.25, "Navigation state sₜ", size=9.6, bold=True, color="8C6C0B")
    add_text(slide, 0.40, 3.42, 1.68, 0.19, "goal and UAV motion (8-D)", size=7.9, color="muted")

    # Encoders and fusion.
    add_encoder(slide, 2.66, 0.70, 2.03, 1.02, "blue", "blue_fill", "LiDAR CNN", "128-D")
    add_encoder(slide, 2.66, 1.91, 2.03, 1.02, "orange", "orange_fill", "Dynamic-obstacle MLP", "64-D")
    add_line(slide, 2.18, 1.15, 2.66, 1.15, color="blue", width=1.1, arrow=True)
    add_line(slide, 2.18, 2.26, 2.66, 2.42, color="orange", width=1.1, arrow=True)

    # Lines are intentionally orthogonal and kept behind the main modules.
    add_route(slide, [(4.69, 1.21), (4.92, 1.21), (4.92, 2.43), (5.07, 2.43)], color="blue", width=1.0)
    add_line(slide, 4.69, 2.42, 5.07, 2.43, color="orange", width=1.0, arrow=True)
    add_route(slide, [(2.18, 3.37), (4.80, 3.37), (4.80, 2.55), (5.07, 2.55)], color="gold", width=1.0)
    concat = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.98), Inches(2.34), Inches(0.30), Inches(0.30))
    set_shape_style(concat, "navy", "white", width=1.1)
    set_text(concat, "⊕", size=13.0, color="navy")

    add_box(slide, 5.43, 2.05, 1.44, 0.86, "navy", "navy_fill", width=1.15)
    add_text(slide, 5.55, 2.20, 1.20, 0.25, "Fusion MLP", size=10.3, bold=True, color="navy")
    add_text(slide, 5.55, 2.51, 1.20, 0.20, "feature fusion", size=8.2, color="muted")
    add_line(slide, 5.28, 2.49, 5.43, 2.49, color="navy", width=1.0, arrow=True)

    add_box(slide, 7.08, 2.05, 0.85, 0.86, "navy", "white", width=1.15)
    add_text(slide, 7.16, 2.20, 0.69, 0.26, "zₜ", size=13.0, bold=True, italic=True, color="navy")
    add_text(slide, 7.16, 2.52, 0.69, 0.20, "256-D", size=8.3, color="muted")
    add_line(slide, 6.87, 2.49, 7.08, 2.49, color="navy", width=1.0, arrow=True)

    # Shared representation fan-out.
    add_line(slide, 7.93, 2.49, 8.20, 2.49, color="line", width=1.0)
    add_line(slide, 8.20, 1.21, 8.20, 3.90, color="line", width=1.0)
    add_network(slide, 8.44, 0.74, 2.12, 0.98, "navy", "navy_fill", "Beta policy actor", "αₜ, βₜ")
    add_network(slide, 8.44, 2.23, 2.12, 0.82, "green", "green_fill", "Observation-conditioned critic", "Vψ(oₜ)")
    add_line(slide, 8.20, 1.23, 8.44, 1.23, color="navy", width=1.0, arrow=True)
    add_line(slide, 8.20, 2.64, 8.44, 2.64, color="green", width=1.0, arrow=True)

    add_box(slide, 8.44, 3.42, 2.12, 1.03, "teal", "teal_fill", dashed=True, width=1.15)
    add_text(slide, 8.58, 3.53, 1.84, 0.38, "Multi-horizon outcome\npredictor", size=9.4, bold=True, color="teal")
    add_text(slide, 8.58, 3.94, 1.84, 0.19, "h ∈ {1, 3, 5}", size=8.4)
    add_text(slide, 8.52, 4.16, 1.96, 0.18, "collision · stuck · clearance · progress", size=7.2, color="muted")
    add_line(slide, 8.20, 3.92, 8.44, 3.92, color="teal", width=1.0, arrow=True)

    # Action and action-conditioning branch.
    add_box(slide, 11.36, 0.82, 1.78, 0.88, "navy", "white", width=1.15)
    add_text(slide, 11.47, 0.95, 1.56, 0.23, "Bounded 3-D velocity", size=9.7, bold=True, color="navy")
    add_text(slide, 11.45, 1.23, 1.60, 0.20, "aₜᴳ = vlim(2ãₜ − 1)", size=8.6)
    add_text(slide, 11.45, 1.46, 1.60, 0.16, "vlim = 2 m/s", size=7.6, color="muted")
    add_line(slide, 10.56, 1.23, 11.36, 1.23, color="navy", width=1.15, arrow=True)
    add_route(slide, [(10.90, 1.23), (10.90, 3.20), (10.70, 3.20), (10.70, 3.92), (10.56, 3.92)], color="teal", width=0.9, dashed=True)
    add_text(slide, 10.78, 3.08, 1.05, 0.19, "current ãₜ", size=7.4, color="teal", align=PP_ALIGN.LEFT)

    # Training-only supervision band.
    add_box(slide, 0.12, 4.93, 13.05, 1.90, "band_edge", "band", dashed=True, width=1.0)
    add_text(slide, 0.34, 5.03, 2.25, 0.25, "Training-only supervision", size=11.1, bold=True, color="coral", align=PP_ALIGN.LEFT)

    add_box(slide, 3.30, 5.54, 2.10, 0.94, "line", "white", width=1.1)
    add_text(slide, 3.44, 5.66, 1.82, 0.24, "Joint optimization", size=10.2, bold=True)
    add_text(slide, 3.43, 5.96, 1.84, 0.22, "L = LPPO + λMH LMH", size=9.1, italic=True)
    add_text(slide, 3.43, 6.22, 1.84, 0.18, "shared predictive representation", size=7.3, color="muted")

    add_box(slide, 6.04, 5.33, 1.82, 0.61, "navy", "navy_fill", width=1.0)
    add_text(slide, 6.15, 5.43, 1.60, 0.22, "PPO loss LPPO", size=9.1, bold=True, color="navy")
    add_text(slide, 6.15, 5.69, 1.60, 0.16, "policy and value learning", size=7.0, color="muted")

    add_box(slide, 6.04, 6.10, 1.82, 0.61, "teal", "teal_fill", width=1.0)
    add_text(slide, 6.15, 6.20, 1.60, 0.22, "Outcome loss LMH", size=8.9, bold=True, color="teal")
    add_text(slide, 6.15, 6.46, 1.60, 0.16, "predictive supervision", size=7.0, color="muted")

    add_box(slide, 8.39, 5.33, 1.95, 0.61, "coral", "coral_fill", width=1.0)
    add_vo_icon(slide, 8.50, 5.43)
    add_text(slide, 9.05, 5.40, 1.20, 0.24, "TTC-aware 3D-VO", size=8.2, bold=True, color="coral")
    add_text(slide, 9.05, 5.68, 1.20, 0.17, "risk reward", size=7.6, bold=True, color="coral")

    add_box(slide, 8.39, 6.10, 1.95, 0.61, "teal", "teal_fill", width=1.0)
    add_text(slide, 8.51, 6.20, 1.71, 0.22, "Multi-horizon targets", size=8.8, bold=True, color="teal")
    add_text(slide, 8.51, 6.46, 1.71, 0.16, "from rollout sequences", size=7.0, color="muted")

    add_box(slide, 10.92, 5.43, 2.08, 1.20, "line", "white", width=1.0)
    add_text(slide, 11.05, 5.55, 1.82, 0.25, "Parallel curriculum rollouts", size=9.2, bold=True)
    add_text(slide, 11.05, 5.84, 1.82, 0.19, "dynamic + non-convex obstacles", size=7.2, color="muted")
    # Minimal editable environment glyph.
    set_shape_style(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.43), Inches(6.13), Inches(0.15), Inches(0.30)), "navy", "blue", width=0.5)
    set_shape_style(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.76), Inches(6.26), Inches(0.42), Inches(0.15)), "coral", "orange", width=0.5)
    add_line(slide, 12.42, 6.43, 12.42, 6.14, color="gold", width=1.5)
    add_line(slide, 12.42, 6.14, 12.57, 6.14, color="gold", width=1.5)
    add_line(slide, 12.57, 6.14, 12.57, 6.31, color="gold", width=1.5)
    add_line(slide, 12.57, 6.31, 12.74, 6.31, color="gold", width=1.5)

    # Two clean supervision chains, right to left.
    add_line(slide, 10.92, 5.73, 10.34, 5.64, color="coral", width=1.0, arrow=True)
    add_line(slide, 10.92, 6.27, 10.34, 6.40, color="teal", width=1.0, arrow=True)
    add_line(slide, 8.39, 5.64, 7.86, 5.64, color="coral", width=1.0, arrow=True)
    add_line(slide, 8.39, 6.40, 7.86, 6.40, color="teal", width=1.0, arrow=True)
    add_line(slide, 6.04, 5.64, 5.40, 5.82, color="navy", width=1.0, arrow=True)
    add_line(slide, 6.04, 6.40, 5.40, 6.18, color="teal", width=1.0, arrow=True)

    # Joint optimization updates the common representation. Exact per-head
    # gradient routing remains documented in the Methodology text.
    add_route(slide, [(4.35, 5.54), (4.35, 4.70), (7.50, 4.70), (7.50, 2.91)], color="navy", width=0.9, dashed=True)
    add_text(slide, 5.15, 4.54, 1.85, 0.18, "joint representation update", size=7.2, color="navy")

    # Policy action closes the rollout loop.
    add_route(slide, [(12.85, 1.70), (12.85, 4.76), (12.35, 4.76), (12.35, 5.43)], color="line", width=0.9)
    add_text(slide, 12.55, 4.74, 0.46, 0.18, "action", size=7.0, color="muted")

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
